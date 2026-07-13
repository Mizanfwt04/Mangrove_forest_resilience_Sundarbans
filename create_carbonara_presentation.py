#!/usr/bin/env python3
"""Generate CarbonARA RA interview presentation for Md Mizanur Rahman."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUTPUT = "/workspace/CarbonARA_Interview_MdMizanRahman.pptx"

# Brand colours (forest green + dark slate)
GREEN = RGBColor(0x1B, 0x5E, 0x20)
DARK = RGBColor(0x26, 0x32, 0x38)
ACCENT = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x55, 0x55, 0x55)


def set_title_style(shape, size=32):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = GREEN


def add_bullets(text_frame, items, size=18, level0=True):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0 if level0 else 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_speaker_notes(slide, notes):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title, 36)
    sub = slide.placeholders[1]
    sub.text = subtitle
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
    return slide


def add_content_slide(prs, title, bullets, notes=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title, 28)
    body = slide.placeholders[1].text_frame
    all_items = ([subtitle] if subtitle else []) + bullets
    add_bullets(body, all_items, size=17)
    if subtitle and body.paragraphs:
        body.paragraphs[0].font.italic = True
        body.paragraphs[0].font.color.rgb = ACCENT
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tx.text_frame.text = title
    set_title_style(tx, 28)

    # Left column
    lx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.5))
    lf = lx.text_frame
    lf.text = left_title
    lf.paragraphs[0].font.bold = True
    lf.paragraphs[0].font.size = Pt(20)
    lf.paragraphs[0].font.color.rgb = ACCENT
    for item in left_items:
        p = lf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    # Right column
    rx = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(5.5))
    rf = rx.text_frame
    rf.text = right_title
    rf.paragraphs[0].font.bold = True
    rf.paragraphs[0].font.size = Pt(20)
    rf.paragraphs[0].font.color.rgb = ACCENT
    for item in right_items:
        p = rf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    if notes:
        add_speaker_notes(slide, notes)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = add_title_slide(
        prs,
        "CarbonARA Research Associate Interview",
        "Md Mizanur Rahman, PhD\n"
        "When tropical forests look stable but stop sequestering carbon\n"
        "King's College London  |  July 2026",
    )
    add_speaker_notes(
        s1,
        "~30 seconds. Thank the panel. One-line hook: I study hidden loss of carbon "
        "sink strength in tropical forests using long-term field data fused with Earth observation.",
    )

    # Slide 2
    add_content_slide(
        prs,
        "Why this matters: Amazon & the Global Stocktake",
        [
            "CEOS identifies the Amazon as a potential near-term tipping point for terrestrial carbon emissions",
            "Eastern Amazon (Pará, Santarém): evidence of weakening sink and dry-season carbon dynamics",
            "Large uncertainty remains in CO₂, CH₄, and N₂O fluxes across land covers and fire regimes",
            "CarbonARA: integrated ground, airborne, and satellite campaign to close this gap",
            "My research directly targets: when and where forests lose carbon sink function before cover changes",
        ],
        notes="~1 min. Show you know CarbonARA rationale. Mention Martin Wooster / ESA / INPE / UFOPA collaboration briefly.",
    )

    # Slide 3
    add_content_slide(
        prs,
        "My background & relevant experience",
        [
            "PhD Forest & Biomaterials Science, Kyoto University (2019)",
            "15+ years tropical forest carbon, biodiversity, and remote sensing",
            "Research Fellow, PolyU (2024–present): mangrove resilience, UAV/TLS, GeoAI",
            "Postdoctoral work: HKU (RGC Fellow), HKUST Guangzhou — carbon accounting & ML",
            "NASA GSFC Visiting Fellow (2016); collaborations with US Forest Service, Bangladesh Forest Dept.",
            "20+ peer-reviewed papers | h-index 14 | 1,450+ citations",
            "My role: lead analyst, field campaign leader, pipeline developer, lead/corresponding author",
        ],
        notes="~1 min. Emphasise what YOU did, not only team achievements.",
    )

    # Slide 4
    add_content_slide(
        prs,
        "Long-term field carbon programme — Sundarbans",
        [
            "Quantified above- and below-ground carbon pools across salinity and vegetation zones",
            "Led / coordinated field inventory: tree DBH, species, soil cores, traits",
            "Repeated measurements reveal carbon dynamics beyond snapshot stock estimates",
            "Key finding — Hidden sink degradation:",
            "   • ~23% of plots shifted to net carbon sources",
            "   • Persistent forest cover — sink loss not visible from cover change alone",
            "Structural equation models: hydro-climatic stress → species/trait change → carbon loss",
            "My role: designed analysis, integrated 68 GEE hydro-climatic indices, led SEM & interpretation",
        ],
        subtitle="13-year | 150 permanent plots | world's largest contiguous mangrove forest",
        notes="~2 min. This is your strongest CarbonARA-relevant result. Stress mechanistic understanding.",
    )

    # Slide 5
    add_content_slide(
        prs,
        "Resilience & stability from satellite time series",
        [
            "Framework: MODIS kNDVI (2000–2024) → STL decomposition → perturbation detection",
            "Recovery dynamics: exponential fit, perturbation frequency, disturbance indices",
            "Critical slowing down: sliding-window AC1 (λ) and variance — early warning signals",
            "Published: Rahman et al. 2026, Communications Earth & Environment",
            "   • Functional composition (canopy height, SLA) strongest resilience drivers",
            "   • ~10–15% of Sundarbans area shows declining resilience",
            "My role: coded full pipeline (Python/R), lead author, SEM framework, spatial mapping",
            "CarbonARA link: same logic applicable to SIF, VOD, thermal + tower flux validation",
        ],
        notes="~1.5 min. Point to figure placeholder — insert kNDVI map or AC1 map from your notebook.",
    )

    # Slide 6
    add_content_slide(
        prs,
        "Tropical Forest Carbon Stability — including Amazon",
        [
            "Project: Tropical_forest_carbon_stability (Sundarbans + Amazon)",
            "Amazon: applied stability/resilience metrics to tropical terra firme forest landscapes",
            "Same workflow: long-term vegetation time series → perturbation → AC1/λ → stability class",
            "Compares regions losing recovery capacity despite structurally intact canopy",
            "Different stressors (cyclones vs drought/fire) — shared mechanism: carbon-climate feedback risk",
            "My role: extended pipeline to Amazon AOI, processed satellite stacks, spatial interpretation",
            "[INSERT FIGURE: Amazon stability map + study area inset]",
            "Direct relevance: eastern Amazon transition zone targeted by CarbonARA (Santarém, Pará)",
        ],
        subtitle="Multi-biome extension of the same reproducible pipeline",
        notes="~1.5 min. Replace bracketed figure with your actual Amazon map. Be honest: satellite pilot, not Amazon fieldwork yet.",
    )

    # Slide 7
    add_two_column_slide(
        prs,
        "Multi-sensor Earth observation & validation experience",
        "Data & methods I use",
        [
            "Optical: Landsat 5–8, Sentinel-2, MODIS, PlanetScope, WorldView-2",
            "Radar: Sentinel-1, TanDEM-X; canopy height from GEDI",
            "Platforms: Google Earth Engine, Python, R, ENVI, SNAP",
            "Analytics: ML (XGBoost, RF), SEM (lavaan), time-series decomposition",
            "Proximal: UAV (HK certified), terrestrial LiDAR, hyperspectral",
            "Global product experience: mangrove soil carbon map (Sanderman et al. 2018)",
        ],
        "CarbonARA measurements I can work with",
        [
            "Tower fluxes: CO₂, CH₄ (LICOR), energy balance (IRGASON)",
            "Vegetation: FLoX (SIF), L-band radiometer & GNSS (VOD)",
            "Thermal: KT15 LST radiometers",
            "Airborne: BAS Twin Otter — SIF, methane hyperspectral, in-situ GHGs",
            "Satellites: Sentinel-2/3/5P, FLEX, BIOMASS, MicroCarb, MetOp/IASI",
            "Fire: mobile roving sensors for smoke & particulates",
        ],
        notes="~1 min. Position yourself as ready to validate satellite products — core RA task.",
    )

    # Slide 8
    add_content_slide(
        prs,
        "Selected publications",
        [
            "Rahman et al. (2026) Comms Earth & Environ. — mangrove resilience, kNDVI, SEM",
            "Rahman et al. (2024) Global Change Biology — functional composition & ecosystem function",
            "Rahman et al. (2021) Nature Communications — co-benefits of mangrove protection",
            "Rahman et al. (2019) Remote Sens. Ecol. Conserv. — high-res carbon mapping (WV-2, TanDEM-X)",
            "Sanderman et al. (2018) Env. Res. Lett. — global mangrove soil carbon (co-author)",
            "In preparation: multidecadal carbon mapping; dominant species & climate effects on sequestration",
        ],
        notes="~30 sec. Don't read the list — highlight 2–3 most relevant to carbon stability and EO.",
    )

    # Slide 9
    add_two_column_slide(
        prs,
        "Where I fit in CarbonARA",
        "What I bring",
        [
            "Tropical forest carbon ecology & long-term inventory",
            "EO time-series, resilience metrics, ML & SEM",
            "Reproducible geospatial pipelines (GEE, Python, R)",
            "Field + proximal sensing (inventory, drone, TLS)",
            "International collaboration & scientific writing",
            "Willing to travel to Brazil; available from July 2026",
        ],
        "CarbonARA science themes",
        [
            "(I) Landscape-scale (non-fire) GHG fluxes",
            "(II) Vegetation processes & carbon storage",
            "(III) Fire magnitude, emissions & impact",
            "(IV) Satellite product calibration & validation",
            "Study area: 100×100 km Santarém, Pará",
            "Sites: LBA Tapajós KM67 (primary) + UFOPA (secondary)",
        ],
        notes="~1 min. Show you've read the CarbonARA website. Mention learning from team on eddy covariance & fire science.",
    )

    # Slide 10 — Science Q1
    add_content_slide(
        prs,
        "Science Question 1",
        [
            "Gap: Forests can remain structurally intact while losing sink strength — hard to detect from cover or biomass maps alone",
            "Importance: eastern Amazon may be approaching sink-to-source transition; early warning supports climate policy & modelling",
            "CarbonARA data to address this:",
            "   • FLoX (SIF) + L-band radiometer & GNSS (VOD) + KT15 (LST) at tower sites",
            "   • Continuous CO₂ fluxes at Tapajós KM67 (primary) vs UFOPA (secondary)",
            "   • Airborne SIF/thermal + Sentinel-3 SLSTR for landscape scaling",
            "Approach: test whether physiological/structural stress signals precede flux reversal — building on my hidden sink degradation & AC1 framework",
        ],
        subtitle="Can integrated SIF, VOD, and thermal observations detect sink degradation before flux towers record net carbon loss?",
        notes="~1.5 min. Required slide per interview instructions. Justify why important.",
    )

    # Slide 11 — Science Q2
    add_content_slide(
        prs,
        "Science Question 2",
        [
            "Gap: Single flux towers cannot represent heterogeneous 100×100 km mosaics of intact, degraded, and agricultural land",
            "Importance: CH₄ and CO₂ budgets differ strongly by land cover; upscaling errors bias Global Stocktake assessments",
            "CarbonARA data to address this:",
            "   • Paired towers: primary (Tapajós) vs secondary/degraded (UFOPA) — CO₂, CH₄, energy balance",
            "   • BAS Twin Otter in-situ GHG sampling + Telops methane hyperspectral",
            "   • Sentinel-5P TROPOMI (CH₄, CO) & MicroCarb (CO₂) for validation",
            "   • Sentinel-2 land-cover stratification within the AOO",
            "Approach: stratified upscaling model linking my carbon-stability classes to flux observations",
        ],
        subtitle="How do primary vs secondary forest differ in landscape-scale CO₂ and CH₄ exchange, and can airborne–tower fusion reduce satellite upscaling errors?",
        notes="~1.5 min. Shows you understand scale mismatch problem central to CarbonARA design.",
    )

    # Slide 12
    add_content_slide(
        prs,
        "What I would contribute as Research Associate",
        [
            "Process & QC harmonised tower, airborne, and satellite datasets for the Santarém domain",
            "Lead analyses linking vegetation structure, SIF/VOD/thermal signals, and carbon flux dynamics",
            "Extend my tropical carbon-stability framework using CarbonARA's unique co-located measurements",
            "Support satellite validation (FLEX, BIOMASS, S5P, MicroCarb) and manuscript preparation",
            "Collaborate with INPE, UFOPA, BAS, and NCEO teams; contribute to ESA deliverables",
            "Available from 1 July 2026 | committed to fieldwork in Brazil & researcher development at King's",
            "Thank you — I welcome your questions",
        ],
        notes="~1 min. End confidently. Prepare 2–3 questions for the panel.",
    )

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
